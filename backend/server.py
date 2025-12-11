"""
URS Backend Server - Python FastAPI server with Tavily integration and AI skills (Gemini/OpenAI)
"""

from fastapi import FastAPI, HTTPException, Request, File, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse
from pydantic import BaseModel
from typing import Optional, List, Dict, Any
import os
import asyncio
import aiohttp
from datetime import datetime
import json
import logging
from dotenv import load_dotenv
import io
import tempfile

# Tavily API imports
try:
    from tavily import TavilyClient
except ImportError:
    print("Tavily not installed. Install with: pip install tavily-python")
    TavilyClient = None

# Google Gemini imports
try:
    import google.generativeai as genai
except ImportError:
    print("Google Generative AI not installed. Install with: pip install google-generativeai")
    genai = None

# OpenAI imports (fallback)
try:
    from openai import OpenAI
    openai_available = True
except ImportError:
    OpenAI = None
    openai_available = False

# File processing imports
try:
    from docx import Document as DocxDocument
except ImportError:
    DocxDocument = None

try:
    from PyPDF2 import PdfReader
except ImportError:
    PdfReader = None

try:
    from openpyxl import load_workbook
except ImportError:
    load_workbook = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    import pandas as pd
except ImportError:
    pd = None

# Load environment variables
load_dotenv()

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Import Tavily MCP Client
try:
    from tavily_mcp_client import TavilyMCPClient
except ImportError:
    logger.warning("Tavily MCP Client not available")
    TavilyMCPClient = None

# Initialize FastAPI app
app = FastAPI(title="URS Backend API", version="1.0.0")

# Configure CORS
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# API Configuration
TAVILY_API_KEY = os.getenv("TAVILY_API_KEY", "")
API_KEY = os.getenv("API_KEY", "") or os.getenv("GOOGLE_API_KEY", "")
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY", "")
MCP_REMOTE_URL = os.getenv("MCP_REMOTE_URL", "")

# File Upload Configuration
MAX_FILE_SIZE_MB = float(os.getenv("MAX_FILE_SIZE_MB", "5"))
MAX_CONTENT_LENGTH = int(os.getenv("MAX_CONTENT_LENGTH", "50000"))

# AI Analysis Configuration
MAX_ANALYSIS_LENGTH = int(os.getenv("MAX_ANALYSIS_LENGTH", "10000"))

# Initialize Clients
tavily_client = TavilyClient(api_key=TAVILY_API_KEY) if TavilyClient and TAVILY_API_KEY else None

# Google API Configuration
google_api_available = bool(API_KEY)
if google_api_available:
    logger.info("Google API key configured for REST API calls")

# Initialize OpenAI (Fallback)
openai_client = None
if openai_available and OPENAI_API_KEY:
    try:
        openai_client = OpenAI(api_key=OPENAI_API_KEY)
        logger.info("OpenAI client initialized successfully")
    except Exception as e:
        logger.error(f"Failed to initialize OpenAI: {e}")

# Initialize Tavily Expert MCP Client
tavily_expert_client = None
if TavilyMCPClient and MCP_REMOTE_URL:
    try:
        tavily_expert_client = TavilyMCPClient(MCP_REMOTE_URL)
        logger.info(f"Tavily Expert MCP Client initialized with URL: {MCP_REMOTE_URL}")
    except Exception as e:
        logger.error(f"Failed to initialize Tavily Expert MCP: {e}")

# --- Helper Functions for AI ---

async def generate_ai_response(prompt: str, system_prompt: str = "") -> str:
    """Generate response using available AI model (Google REST API first, then OpenAI)"""
    
    # Try Google Generative AI REST API
    if API_KEY:
        try:
            full_prompt = f"{system_prompt}\n\n{prompt}" if system_prompt else prompt
            
            # Use REST API endpoint with key parameter
            url = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-pro-latest:generateContent?key={API_KEY}"
            
            payload = {
                "contents": [{
                    "parts": [{
                        "text": full_prompt
                    }]
                }],
                "generationConfig": {
                    "temperature": 0.7,
                    "topP": 0.95,
                    "maxOutputTokens": 2048,
                }
            }
            
            async with aiohttp.ClientSession() as session:
                async with session.post(url, json=payload) as response:
                    if response.status == 200:
                        data = await response.json()
                        if data.get("candidates") and len(data["candidates"]) > 0:
                            text = data["candidates"][0]["content"]["parts"][0]["text"]
                            return text
                        else:
                            logger.warning("Google API returned empty candidates")
                            raise Exception("Empty response from Google API")
                    else:
                        error_text = await response.text()
                        logger.error(f"Google API error: {response.status} - {error_text}")
                        raise Exception(f"Google API returned {response.status}")
                        
        except Exception as e:
            logger.error(f"Google API generation error: {type(e).__name__}: {str(e)}")
            # Fall through to OpenAI if available
    
    # Try OpenAI
    if openai_client:
        try:
            messages = []
            if system_prompt:
                messages.append({"role": "system", "content": system_prompt})
            messages.append({"role": "user", "content": prompt})
            
            response = openai_client.chat.completions.create(
                model="gpt-3.5-turbo",
                messages=messages,
                max_tokens=2048,
                temperature=0.7
            )
            return response.choices[0].message.content
        except Exception as e:
            logger.error(f"OpenAI generation error: {type(e).__name__}: {str(e)}")
            
    return "AI service unavailable. Please check your API keys."

# --- Request Models ---

class TavilySearchRequest(BaseModel):
    query: str
    search_depth: Optional[str] = "basic"
    include_answer: Optional[bool] = True
    include_images: Optional[bool] = False
    max_results: Optional[int] = 5
    context_url: Optional[str] = None

class TavilyExtractRequest(BaseModel):
    url: str
    extract_depth: Optional[str] = "basic"
    include_images: Optional[bool] = False

class SummarizeRequest(BaseModel):
    content: str
    url: Optional[str] = None
    use_page_content: Optional[bool] = False
    max_length: Optional[int] = 500

class FactCheckRequest(BaseModel):
    statement: str
    context_url: Optional[str] = None

class CreateReportRequest(BaseModel):
    content: str
    url: Optional[str] = None
    report_type: Optional[str] = "comprehensive"

class SkillRequest(BaseModel):
    input: str
    skill: str
    context: Optional[Dict[str, Any]] = None

# --- Endpoints ---

@app.get("/health")
async def health_check():
    return {
        "status": "healthy",
        "timestamp": datetime.now().isoformat(),
        "services": {
            "tavily": "available" if tavily_client else "unavailable",
            "google_api": "available" if API_KEY else "unavailable",
            "openai": "available" if openai_client else "unavailable",
            "tavily_expert_mcp": "available" if tavily_expert_client else "unavailable"
        },
        "file_upload_limits": {
            "max_file_size_mb": MAX_FILE_SIZE_MB,
            "max_content_length": MAX_CONTENT_LENGTH
        },
        "mcp_remote_url": MCP_REMOTE_URL if MCP_REMOTE_URL else None
    }

@app.post("/api/tavily/search")
async def tavily_search(request: TavilySearchRequest):
    if not tavily_client:
        raise HTTPException(status_code=503, detail="Tavily client not initialized")
    
    try:
        results = tavily_client.search(
            query=request.query,
            search_depth=request.search_depth,
            include_answer=request.include_answer,
            include_images=request.include_images,
            max_results=request.max_results
        )
        return {
            "success": True,
            "query": request.query,
            "results": results.get("results", []),
            "answer": results.get("answer", None),
            "images": results.get("images", []) if request.include_images else []
        }
    except Exception as e:
        logger.error(f"Tavily search error: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/tavily/extract")
async def tavily_extract(request: TavilyExtractRequest):
    if not tavily_client:
        raise HTTPException(status_code=503, detail="Tavily client not initialized")
    
    try:
        results = tavily_client.extract(
            urls=[request.url],
            include_images=request.include_images
        )
        return {
            "success": True,
            "url": request.url,
            "content": results.get("results", [{}])[0] if results.get("results") else {},
            "images": results.get("images", []) if request.include_images else []
        }
    except Exception as e:
        logger.error(f"Tavily extract error: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/skills/summarize")
async def summarize_content(request: SummarizeRequest):
    try:
        content = request.content
        
        # Extract content from URL if needed
        if request.use_page_content and request.url and tavily_client:
            extract_result = tavily_client.extract(urls=[request.url])
            if extract_result.get("results"):
                content = extract_result["results"][0].get("content", content)
        
        system_prompt = "You are a helpful assistant that creates concise summaries."
        prompt = f"Please summarize the following content in {request.max_length} words or less:\n\n{content}"
        
        summary = await generate_ai_response(prompt, system_prompt)
        
        return {
            "success": True,
            "summary": summary,
            "original_length": len(content),
            "summary_length": len(summary)
        }
    except Exception as e:
        logger.error(f"Summarize error: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/skills/fact_check")
async def fact_check(request: FactCheckRequest):
    try:
        # 1. Search with Tavily
        if not tavily_client:
            raise HTTPException(status_code=503, detail="Tavily client not initialized")
            
        # Generate a search query if the statement is too long
        search_query = request.statement
        if len(search_query) > 300:
             if google_api_available:
                query_prompt = f"Generate a short search query (max 30 words) to fact check this statement: {request.statement[:500]}..."
                search_query = await generate_ai_response(query_prompt)
             else:
                search_query = request.statement[:300]

        search_results = tavily_client.search(
            query=f"fact check {search_query}",
            search_depth="advanced",
            include_answer=True,
            max_results=5
        )
        
        context = search_results.get('answer', '')
        results_list = search_results.get('results', [])
        
        # 2. Analyze with AI
        system_prompt = "You are a professional fact-checker. Analyze the statement against the provided search results. Be objective."
        prompt = f"""
        Statement to verify: "{request.statement}"
        
        Search Results Context:
        {context}
        
        Detailed Sources:
        {json.dumps(results_list[:3], indent=2)}
        
        Please provide:
        1. A verdict (True, False, Partially True, Unverified)
        2. A detailed explanation citing the sources.
        """
        
        analysis = await generate_ai_response(prompt, system_prompt)
        
        return {
            "success": True,
            "statement": request.statement,
            "fact_check": {
                "verdict": "See analysis", # Simpler to let the text explain
                "explanation": analysis,
                "sources": results_list[:3]
            }
        }
    except Exception as e:
        logger.error(f"Fact check error: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/skills/create_report")
async def create_report(request: CreateReportRequest):
    try:
        content = request.content
        
        # 1. Gather context
        if request.url and tavily_client:
            extract_result = tavily_client.extract(urls=[request.url])
            if extract_result.get("results"):
                content += f"\n\nSource Content from {request.url}:\n{extract_result['results'][0].get('content', '')}"
        
        # 2. Search for broader context
        search_context = ""
        if tavily_client:
            # Generate a search query from the content (using AI if possible, or simple truncation)
            query_prompt = f"Generate a good search query to find more info about this topic: {content[:100]}..."
            search_query = await generate_ai_response(query_prompt) if google_api_available else content[:100]
            
            search_results = tavily_client.search(
                query=search_query,
                search_depth="advanced",
                max_results=5
            )
            search_context = json.dumps(search_results.get("results", []), indent=2)

        # 3. Generate Report
        system_prompt = "You are a professional research analyst. Create well-structured, comprehensive reports."
        prompt = f"""
        Create a {request.report_type} report based on the following:
        
        Primary Content:
        {content[:10000]} # Truncate to avoid limit issues if massive
        
        Additional Search Context:
        {search_context}
        
        Structure the report with:
        - Executive Summary
        - Key Findings
        - Detailed Analysis
        - Implications/Recommendations
        - Conclusion
        """
        
        report = await generate_ai_response(prompt, system_prompt)
        
        return {
            "success": True,
            "report": report,
            "timestamp": datetime.now().isoformat()
        }
    except Exception as e:
        logger.error(f"Report error: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/skills/analyze")
async def analyze_content(request: SkillRequest):
    """Analyze content and provide insights and recommendations"""
    try:
        input_text = request.input
        
        # Truncate input if too long (AI models have token limits)
        # MAX_ANALYSIS_LENGTH is now configured in .env
        original_length = len(input_text)
        
        if len(input_text) > MAX_ANALYSIS_LENGTH:
            logger.info(f"Truncating input from {original_length} to {MAX_ANALYSIS_LENGTH} chars for analysis")
            input_text = input_text[:MAX_ANALYSIS_LENGTH]
            truncated = True
        else:
            truncated = False
        
        # Use AI model for analysis
        system_prompt = "You are an analytical assistant. Provide clear, concise insights and actionable recommendations."
        
        prompt = f"""Analyze the following content and provide:

1. **Summary** (2-3 sentences)
2. **Key Insights** (3-5 main points)
3. **Recommendations** (actionable suggestions)

Content ({len(input_text)} characters):
{input_text}
"""
        
        if truncated:
            prompt += f"\n\n[Note: Content was truncated from {original_length:,} to {MAX_ANALYSIS_LENGTH:,} characters for analysis]"
        
        try:
            analysis = await generate_ai_response(prompt, system_prompt)
        except Exception as ai_error:
            logger.error(f"AI generation error in analyze: {ai_error}")
            # Fallback to basic analysis
            word_count = len(input_text.split())
            analysis = f"""## Analysis of Content

**Summary**: The content contains approximately {word_count:,} words across {original_length:,} characters.

**Key Insights**:
- Document type: Text content
- Length: {original_length:,} characters
- Word count: ~{word_count:,} words
{f'- Note: Content was truncated for analysis' if truncated else ''}

**Recommendations**: 
- For more detailed analysis, consider using the 'Summarize' skill first to condense the content
- Break down large documents into smaller sections for focused analysis
- Use specific questions or prompts for targeted insights
"""
        
        # Optional: Use Tavily for additional context (only for shorter inputs)
        search_context = None
        if len(input_text) < 500 and tavily_client:
            try:
                search_results = tavily_client.search(
                    query=input_text[:200],
                    search_depth="basic",
                    include_answer=True,
                    max_results=3
                )
                if search_results.get("answer"):
                    search_context = search_results["answer"]
            except Exception as search_error:
                logger.warning(f"Tavily search error during analysis: {search_error}")
        
        return {
            "success": True,
            "analysis": {
                "result": analysis,
                "original_length": original_length,
                "analyzed_length": len(input_text),
                "truncated": truncated,
                "word_count": len(input_text.split()),
                "search_context": search_context if search_context else None
            }
        }
        
    except Exception as e:
        logger.error(f"Analyze error: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Analysis failed: {str(e)}")

# Placeholder for other tools
@app.post("/api/tavily/map")
async def tavily_map(request: dict):
    return {"message": "Map functionality coming soon"}

@app.post("/api/tavily/crawl")
async def tavily_crawl(request: dict):
    return {"message": "Crawl functionality coming soon"}

# ========== TAVILY EXPERT MCP ENDPOINTS ==========

@app.get("/api/mcp/tools")
async def get_mcp_tools():
    """Get available tools from Tavily Expert MCP server"""
    if not tavily_expert_client:
        raise HTTPException(status_code=503, detail="Tavily Expert MCP not configured")
    
    try:
        tools = await tavily_expert_client.get_available_tools()
        return {
            "success": True,
            "tools": tools,
            "mcp_url": MCP_REMOTE_URL
        }
    except Exception as e:
        logger.error(f"Failed to get MCP tools: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/mcp/call")
async def call_mcp_tool(request: Dict[str, Any]):
    """Call a tool on the Tavily Expert MCP server"""
    if not tavily_expert_client:
        raise HTTPException(status_code=503, detail="Tavily Expert MCP not configured")
    
    tool_name = request.get("tool_name")
    parameters = request.get("parameters", {})
    
    if not tool_name:
        raise HTTPException(status_code=400, detail="tool_name is required")
    
    try:
        result = await tavily_expert_client.call_tool(tool_name, parameters)
        return {
            "success": True,
            "tool_name": tool_name,
            "result": result
        }
    except Exception as e:
        logger.error(f"MCP tool call failed: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/mcp/search")
async def mcp_tavily_search(request: TavilySearchRequest):
    """Enhanced search using Tavily Expert MCP"""
    if not tavily_expert_client:
        # Fallback to regular Tavily
        return await tavily_search(request)
    
    try:
        result = await tavily_expert_client.tavily_search_advanced(
            query=request.query,
            search_depth=request.search_depth,
            include_answer=request.include_answer,
            include_images=request.include_images,
            max_results=request.max_results
        )
        return {
            "success": True,
            "query": request.query,
            "result": result,
            "source": "tavily_expert_mcp"
        }
    except Exception as e:
        logger.error(f"MCP search error: {e}")
        # Fallback to regular Tavily
        return await tavily_search(request)

@app.post("/api/mcp/extract")
async def mcp_tavily_extract(request: TavilyExtractRequest):
    """Enhanced extraction using Tavily Expert MCP"""
    if not tavily_expert_client:
        # Fallback to regular Tavily
        return await tavily_extract(request)
    
    try:
        result = await tavily_expert_client.tavily_extract_advanced(
            urls=[request.url],
            extract_depth=request.extract_depth,
            include_images=request.include_images
        )
        return {
            "success": True,
            "url": request.url,
            "result": result,
            "source": "tavily_expert_mcp"
        }
    except Exception as e:
        logger.error(f"MCP extract error: {e}")
        # Fallback to regular Tavily
        return await tavily_extract(request)

# ========== FILE UPLOAD & PROCESSING ==========

def extract_text_from_txt(file_content: bytes) -> str:
    """Extract text from .txt or .md files"""
    try:
        return file_content.decode('utf-8')
    except UnicodeDecodeError:
        return file_content.decode('latin-1')

def extract_text_from_pdf(file_content: bytes) -> str:
    """Extract text from PDF files"""
    if not PdfReader:
        raise HTTPException(status_code=501, detail="PDF processing not available. Install PyPDF2")
    
    try:
        pdf_file = io.BytesIO(file_content)
        pdf_reader = PdfReader(pdf_file)
        
        text = []
        for page_num, page in enumerate(pdf_reader.pages):
            page_text = page.extract_text()
            if page_text:
                text.append(f"--- Page {page_num + 1} ---\n{page_text}")
        
        return "\n\n".join(text)
    except Exception as e:
        logger.error(f"PDF extraction error: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to extract PDF: {str(e)}")

def extract_text_from_docx(file_content: bytes) -> str:
    """Extract text from .docx files"""
    if not DocxDocument:
        raise HTTPException(status_code=501, detail="DOCX processing not available. Install python-docx")
    
    try:
        docx_file = io.BytesIO(file_content)
        doc = DocxDocument(docx_file)
        
        text = []
        for para in doc.paragraphs:
            if para.text.strip():
                text.append(para.text)
        
        # Also extract text from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = [cell.text for cell in row.cells]
                text.append(" | ".join(row_text))
        
        return "\n\n".join(text)
    except Exception as e:
        logger.error(f"DOCX extraction error: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to extract DOCX: {str(e)}")

def extract_text_from_xlsx(file_content: bytes) -> str:
    """Extract text from .xlsx files"""
    if not load_workbook:
        raise HTTPException(status_code=501, detail="XLSX processing not available. Install openpyxl")
    
    try:
        xlsx_file = io.BytesIO(file_content)
        wb = load_workbook(xlsx_file, data_only=True)
        
        text = []
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            text.append(f"=== Sheet: {sheet_name} ===\n")
            
            for row in sheet.iter_rows(values_only=True):
                row_text = [str(cell) if cell is not None else "" for cell in row]
                if any(row_text):  # Skip empty rows
                    text.append(" | ".join(row_text))
        
        return "\n".join(text)
    except Exception as e:
        logger.error(f"XLSX extraction error: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to extract XLSX: {str(e)}")

def extract_text_from_csv(file_content: bytes) -> str:
    """Extract text from .csv files"""
    if not pd:
        # Fallback to simple text extraction if pandas not available
        return extract_text_from_txt(file_content)
    
    try:
        csv_file = io.BytesIO(file_content)
        df = pd.read_csv(csv_file)
        
        # Convert to string representation
        return df.to_string()
    except Exception as e:
        logger.error(f"CSV extraction error: {e}")
        # Fallback to text extraction
        return extract_text_from_txt(file_content)

def extract_text_from_pptx(file_content: bytes) -> str:
    """Extract text from .pptx files"""
    if not Presentation:
        raise HTTPException(status_code=501, detail="PPTX processing not available. Install python-pptx")
    
    try:
        pptx_file = io.BytesIO(file_content)
        prs = Presentation(pptx_file)
        
        text = []
        for slide_num, slide in enumerate(prs.slides):
            text.append(f"=== Slide {slide_num + 1} ===\n")
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text.append(shape.text)
        
        return "\n\n".join(text)
    except Exception as e:
        logger.error(f"PPTX extraction error: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to extract PPTX: {str(e)}")

@app.post("/api/upload/file")
async def upload_file(file: UploadFile = File(...)):
    """Handle file uploads and extract text content"""
    
    # Use configuration from environment variables
    max_file_size_bytes = int(MAX_FILE_SIZE_MB * 1024 * 1024)  # Convert MB to bytes
    max_content_chars = MAX_CONTENT_LENGTH
    
    file_content = await file.read()
    
    if len(file_content) > max_file_size_bytes:
        raise HTTPException(
            status_code=413,
            detail=f"File too large. Max size is {MAX_FILE_SIZE_MB}MB, got {len(file_content) / 1024 / 1024:.2f}MB"
        )
    
    # Get file extension
    filename = file.filename.lower()
    
    try:
        # Extract text based on file type
        if filename.endswith('.txt') or filename.endswith('.md'):
            extracted_text = extract_text_from_txt(file_content)
        
        elif filename.endswith('.pdf'):
            extracted_text = extract_text_from_pdf(file_content)
        
        elif filename.endswith('.docx'):
            extracted_text = extract_text_from_docx(file_content)
        
        elif filename.endswith('.xlsx'):
            extracted_text = extract_text_from_xlsx(file_content)
        
        elif filename.endswith('.csv'):
            extracted_text = extract_text_from_csv(file_content)
        
        elif filename.endswith('.pptx'):
            extracted_text = extract_text_from_pptx(file_content)
        
        else:
            raise HTTPException(
                status_code=400,
                detail=f"Unsupported file type: {filename.split('.')[-1]}"
            )
        
        # Validate content length after extraction
        content_length = len(extracted_text)
        if content_length > max_content_chars:
            raise HTTPException(
                status_code=413,
                detail=f"Extracted content too long. Max {max_content_chars:,} characters, got {content_length:,} characters. Please use a smaller file or summarize the content first."
            )
        
        return {
            "success": True,
            "filename": file.filename,
            "file_size": len(file_content),
            "file_size_mb": round(len(file_content) / 1024 / 1024, 2),
            "extracted_text": extracted_text,
            "char_count": content_length,
            "content_length": content_length,
            "max_content_length": max_content_chars,
            "max_file_size_mb": MAX_FILE_SIZE_MB,
            "timestamp": datetime.now().isoformat()
        }
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"File processing error: {e}")
        raise HTTPException(status_code=500, detail=f"Error processing file: {str(e)}")

if __name__ == "__main__":
    import uvicorn
    # Use import string to enable reload
    uvicorn.run("server:app", host="0.0.0.0", port=8000, reload=True)